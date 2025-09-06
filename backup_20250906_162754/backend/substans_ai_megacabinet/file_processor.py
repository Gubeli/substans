"""
Système de traitement de fichiers pour substans.ai
Analyse et traite tous types de documents joints aux missions
"""

import os
import json
import mimetypes
from datetime import datetime
from typing import Dict, List, Any, Optional
import hashlib

# Imports pour le traitement de fichiers
try:
    import pandas as pd
    import docx
    from pptx import Presentation
    import PyPDF2
    import markdown
except ImportError as e:
    print(f"Attention: Certaines dépendances manquent pour le traitement de fichiers: {e}")

class FileProcessor:
    """
    Processeur de fichiers pour substans.ai
    Traite tous types de documents et les intègre à la base de connaissances
    """
    
    def __init__(self, knowledge_base_path="/home/ubuntu/substans_ai_megacabinet/knowledge_base"):
        self.knowledge_base_path = knowledge_base_path
        self.processed_files_path = os.path.join(knowledge_base_path, "processed_files")
        self.file_metadata_path = os.path.join(knowledge_base_path, "file_metadata.json")
        
        # Créer les répertoires nécessaires
        os.makedirs(self.processed_files_path, exist_ok=True)
        
        # Types de fichiers supportés
        self.supported_types = {
            'text': ['.txt', '.md', '.markdown'],
            'office': ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'],
            'pdf': ['.pdf'],
            'data': ['.csv', '.json', '.xml']
        }
        
        # Métadonnées des fichiers traités
        self.file_metadata = self.load_file_metadata()
        
    def load_file_metadata(self) -> Dict[str, Any]:
        """Charge les métadonnées des fichiers traités"""
        if os.path.exists(self.file_metadata_path):
            with open(self.file_metadata_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {}
    
    def save_file_metadata(self):
        """Sauvegarde les métadonnées des fichiers"""
        with open(self.file_metadata_path, 'w', encoding='utf-8') as f:
            json.dump(self.file_metadata, f, indent=2, ensure_ascii=False)
    
    def get_file_hash(self, file_path: str) -> str:
        """Calcule le hash d'un fichier pour détecter les doublons"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def detect_file_type(self, file_path: str) -> str:
        """Détecte le type de fichier"""
        _, ext = os.path.splitext(file_path.lower())
        
        for file_type, extensions in self.supported_types.items():
            if ext in extensions:
                return file_type
        
        return 'unknown'
    
    def process_text_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier texte"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Traitement spécial pour Markdown
            if file_path.lower().endswith(('.md', '.markdown')):
                try:
                    html_content = markdown.markdown(content)
                    return {
                        'type': 'markdown',
                        'content': content,
                        'html_content': html_content,
                        'word_count': len(content.split()),
                        'line_count': len(content.split('\n'))
                    }
                except:
                    pass
            
            return {
                'type': 'text',
                'content': content,
                'word_count': len(content.split()),
                'line_count': len(content.split('\n'))
            }
            
        except Exception as e:
            return {'error': f"Erreur lecture fichier texte: {str(e)}"}
    
    def process_pdf_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier PDF"""
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
                
                return {
                    'type': 'pdf',
                    'content': content,
                    'page_count': len(pdf_reader.pages),
                    'word_count': len(content.split()),
                    'metadata': pdf_reader.metadata if hasattr(pdf_reader, 'metadata') else {}
                }
                
        except Exception as e:
            return {'error': f"Erreur lecture PDF: {str(e)}"}
    
    def process_word_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier Word"""
        try:
            doc = docx.Document(file_path)
            
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            # Extraction des tableaux
            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_content.append(table_data)
            
            return {
                'type': 'word',
                'content': content,
                'word_count': len(content.split()),
                'paragraph_count': len(doc.paragraphs),
                'tables': tables_content,
                'table_count': len(tables_content)
            }
            
        except Exception as e:
            return {'error': f"Erreur lecture Word: {str(e)}"}
    
    def process_excel_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier Excel"""
        try:
            # Lire toutes les feuilles
            excel_data = pd.read_excel(file_path, sheet_name=None)
            
            sheets_info = {}
            total_rows = 0
            total_cols = 0
            
            for sheet_name, df in excel_data.items():
                sheets_info[sheet_name] = {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': list(df.columns),
                    'sample_data': df.head(3).to_dict('records') if len(df) > 0 else []
                }
                total_rows += len(df)
                total_cols = max(total_cols, len(df.columns))
            
            return {
                'type': 'excel',
                'sheets': sheets_info,
                'sheet_count': len(excel_data),
                'total_rows': total_rows,
                'max_columns': total_cols
            }
            
        except Exception as e:
            return {'error': f"Erreur lecture Excel: {str(e)}"}
    
    def process_powerpoint_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier PowerPoint"""
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            total_text = ""
            
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                slides_content.append({
                    'slide_number': i + 1,
                    'text': slide_text,
                    'shape_count': len(slide.shapes)
                })
                total_text += slide_text
            
            return {
                'type': 'powerpoint',
                'content': total_text,
                'slides': slides_content,
                'slide_count': len(prs.slides),
                'word_count': len(total_text.split())
            }
            
        except Exception as e:
            return {'error': f"Erreur lecture PowerPoint: {str(e)}"}
    
    def process_csv_file(self, file_path: str) -> Dict[str, Any]:
        """Traite un fichier CSV"""
        try:
            df = pd.read_csv(file_path)
            
            return {
                'type': 'csv',
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': list(df.columns),
                'data_types': df.dtypes.to_dict(),
                'sample_data': df.head(5).to_dict('records'),
                'summary_stats': df.describe().to_dict() if df.select_dtypes(include=['number']).shape[1] > 0 else {}
            }
            
        except Exception as e:
            return {'error': f"Erreur lecture CSV: {str(e)}"}
    
    def process_file(self, file_path: str, mission_id: str = None) -> Dict[str, Any]:
        """Traite un fichier selon son type"""
        
        if not os.path.exists(file_path):
            return {'error': 'Fichier non trouvé'}
        
        # Informations de base
        file_info = {
            'file_path': file_path,
            'file_name': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_hash': self.get_file_hash(file_path),
            'processed_at': datetime.now().isoformat(),
            'mission_id': mission_id
        }
        
        # Vérifier si déjà traité
        if file_info['file_hash'] in self.file_metadata:
            print(f"Fichier déjà traité: {file_info['file_name']}")
            return self.file_metadata[file_info['file_hash']]
        
        # Détecter le type et traiter
        file_type = self.detect_file_type(file_path)
        file_info['file_type'] = file_type
        
        print(f"Traitement du fichier: {file_info['file_name']} (type: {file_type})")
        
        # Traitement selon le type
        if file_type == 'text':
            processing_result = self.process_text_file(file_path)
        elif file_type == 'pdf':
            processing_result = self.process_pdf_file(file_path)
        elif file_type == 'office':
            ext = os.path.splitext(file_path.lower())[1]
            if ext in ['.docx', '.doc']:
                processing_result = self.process_word_file(file_path)
            elif ext in ['.xlsx', '.xls']:
                processing_result = self.process_excel_file(file_path)
            elif ext in ['.pptx', '.ppt']:
                processing_result = self.process_powerpoint_file(file_path)
            else:
                processing_result = {'error': 'Type Office non supporté'}
        elif file_type == 'data':
            if file_path.lower().endswith('.csv'):
                processing_result = self.process_csv_file(file_path)
            else:
                processing_result = {'error': 'Type de données non supporté'}
        else:
            processing_result = {'error': 'Type de fichier non supporté'}
        
        # Combiner les informations
        result = {**file_info, **processing_result}
        
        # Sauvegarder dans les métadonnées
        self.file_metadata[file_info['file_hash']] = result
        self.save_file_metadata()
        
        # Ajouter à la base de connaissances
        self.add_to_knowledge_base(result)
        
        return result
    
    def add_to_knowledge_base(self, file_result: Dict[str, Any]):
        """Ajoute le fichier traité à la base de connaissances"""
        
        if 'error' in file_result:
            return
        
        # Créer un fichier de résumé dans la base de connaissances
        kb_file_path = os.path.join(
            self.processed_files_path, 
            f"{file_result['file_hash']}_summary.json"
        )
        
        # Préparer le résumé pour la base de connaissances
        kb_entry = {
            'file_name': file_result['file_name'],
            'file_type': file_result['file_type'],
            'processed_at': file_result['processed_at'],
            'mission_id': file_result.get('mission_id'),
            'content_summary': self.generate_content_summary(file_result),
            'keywords': self.extract_keywords(file_result),
            'metadata': {
                'file_size': file_result['file_size'],
                'word_count': file_result.get('word_count', 0)
            }
        }
        
        # Sauvegarder
        with open(kb_file_path, 'w', encoding='utf-8') as f:
            json.dump(kb_entry, f, indent=2, ensure_ascii=False)
        
        print(f"Fichier ajouté à la base de connaissances: {kb_file_path}")
    
    def generate_content_summary(self, file_result: Dict[str, Any]) -> str:
        """Génère un résumé du contenu du fichier"""
        
        content = file_result.get('content', '')
        if not content:
            return "Contenu non textuel"
        
        # Résumé simple (premiers 500 caractères)
        summary = content[:500]
        if len(content) > 500:
            summary += "..."
        
        return summary
    
    def extract_keywords(self, file_result: Dict[str, Any]) -> List[str]:
        """Extrait des mots-clés du fichier"""
        
        content = file_result.get('content', '').lower()
        if not content:
            return []
        
        # Mots-clés simples basés sur la fréquence
        words = content.split()
        word_freq = {}
        
        for word in words:
            # Nettoyer le mot
            word = ''.join(c for c in word if c.isalnum())
            if len(word) > 3:  # Ignorer les mots trop courts
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Retourner les 10 mots les plus fréquents
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        return [word for word, freq in sorted_words[:10]]
    
    def get_processed_files_for_mission(self, mission_id: str) -> List[Dict[str, Any]]:
        """Récupère tous les fichiers traités pour une mission"""
        
        mission_files = []
        for file_hash, file_data in self.file_metadata.items():
            if file_data.get('mission_id') == mission_id:
                mission_files.append(file_data)
        
        return mission_files
    
    def search_in_processed_files(self, query: str) -> List[Dict[str, Any]]:
        """Recherche dans les fichiers traités"""
        
        results = []
        query_lower = query.lower()
        
        for file_hash, file_data in self.file_metadata.items():
            # Recherche dans le contenu
            content = file_data.get('content', '').lower()
            if query_lower in content:
                results.append({
                    'file_name': file_data['file_name'],
                    'file_type': file_data['file_type'],
                    'relevance': content.count(query_lower),
                    'file_data': file_data
                })
        
        # Trier par pertinence
        results.sort(key=lambda x: x['relevance'], reverse=True)
        return results

# Test du processeur de fichiers
if __name__ == '__main__':
    processor = FileProcessor()
    
    print("=== Processeur de Fichiers Substans.ai ===")
    print(f"Types supportés: {list(processor.supported_types.keys())}")
    print(f"Base de connaissances: {processor.knowledge_base_path}")
    
    # Test avec un fichier exemple (si disponible)
    test_file = "/home/ubuntu/substans_ai_megacabinet/README.md"
    if os.path.exists(test_file):
        print(f"\nTest avec: {test_file}")
        result = processor.process_file(test_file, mission_id="test_001")
        print(f"Résultat: {result.get('type', 'erreur')}")
        if 'word_count' in result:
            print(f"Mots: {result['word_count']}")
    
    print(f"\nFichiers traités: {len(processor.file_metadata)}")

